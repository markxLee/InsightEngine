#!/usr/bin/env python3
"""Gather content from local files and URLs into combined Markdown.

Usage:
    python3 gather.py --sources sources.json --output collected.md
    python3 gather.py --files doc.pdf report.docx --output collected.md
    python3 gather.py --urls "https://example.com" --output collected.md
    python3 gather.py --files a.pdf --urls "https://x.com" --output collected.md

sources.json format:
    {
        "files": ["path/to/doc.pdf", "path/to/report.docx"],
        "urls": ["https://example.com/page"]
    }

Output: Combined Markdown with source headers, written to --output path.
"""

import argparse
import json
import os
import sys
import time
from pathlib import Path

# ---------------------------------------------------------------------------
# File reading: markitdown first, format-specific fallback second
# ---------------------------------------------------------------------------

def read_file_markitdown(filepath: str) -> str:
    """Read a file using markitdown. Returns extracted text or empty string."""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(filepath)
        return result.text_content or ""
    except Exception:
        return ""


def read_file_fallback(filepath: str) -> str:
    """Format-specific fallback when markitdown fails or returns <100 chars."""
    ext = Path(filepath).suffix.lower()
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(filepath)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            parts = []
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                parts.append(f"### Sheet: {ws.title}\n")
                for row in rows:
                    parts.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
            return "\n".join(parts)
        elif ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                return "\n\n".join(p.extract_text() or "" for p in pdf.pages)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    parts.append(f"### Slide {i}\n" + "\n".join(texts))
            return "\n\n".join(parts)
        elif ext in (".txt", ".md", ".csv", ".tsv"):
            return Path(filepath).read_text(encoding="utf-8")
        else:
            return ""
    except Exception as e:
        return f"[Fallback read error: {e}]"


def read_local_file(filepath: str) -> dict:
    """Read a single local file. Returns {source, content, chars, error}."""
    filepath = os.path.expanduser(filepath)
    if not os.path.isfile(filepath):
        return {"source": filepath, "content": "", "chars": 0, "error": f"File not found: {filepath}"}

    size_mb = os.path.getsize(filepath) / (1024 * 1024)
    if size_mb > 50:
        return {"source": filepath, "content": "", "chars": 0,
                "error": f"File too large ({size_mb:.1f} MB > 50 MB limit). Split or provide specific pages."}

    content = read_file_markitdown(filepath)
    if len(content) < 100:
        content = read_file_fallback(filepath)

    return {"source": filepath, "content": content, "chars": len(content), "error": None}


# ---------------------------------------------------------------------------
# URL fetching: httpx + BeautifulSoup
# ---------------------------------------------------------------------------

def fetch_url(url: str) -> dict:
    """Fetch content from a URL. Returns {source, content, chars, error}."""
    try:
        import httpx
        from bs4 import BeautifulSoup
    except ImportError:
        return {"source": url, "content": "", "chars": 0,
                "error": "Missing httpx or beautifulsoup4. Run: pip install --user httpx beautifulsoup4"}

    try:
        resp = httpx.get(url, timeout=15, follow_redirects=True,
                         headers={"User-Agent": "Mozilla/5.0 (InsightEngine/1.0)"})
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")

        # Remove nav, footer, script, style, cookie banners
        for tag in soup.find_all(["nav", "footer", "script", "style", "noscript", "aside"]):
            tag.decompose()

        # Extract main content
        main = soup.find("main") or soup.find("article") or soup.find("body") or soup
        text = main.get_text(separator="\n", strip=True)

        # Limit to 50,000 chars
        if len(text) > 50_000:
            text = text[:50_000] + "\n\n[... truncated at 50,000 chars]"

        title = soup.title.string.strip() if soup.title and soup.title.string else url
        return {"source": f"{title} ({url})", "content": text, "chars": len(text), "error": None}

    except httpx.TimeoutException:
        return {"source": url, "content": "", "chars": 0, "error": f"Timeout after 15s: {url}"}
    except httpx.HTTPStatusError as e:
        return {"source": url, "content": "", "chars": 0, "error": f"HTTP {e.response.status_code}: {url}"}
    except Exception as e:
        return {"source": url, "content": "", "chars": 0, "error": f"Fetch error: {e}"}


# ---------------------------------------------------------------------------
# Combine and output
# ---------------------------------------------------------------------------

def combine_results(results: list[dict]) -> str:
    """Combine gathered results into structured Markdown."""
    parts = []
    for r in results:
        if r["error"]:
            parts.append(f"## Nguồn: {r['source']}\n> ❌ Lỗi: {r['error']}\n\n---\n")
        else:
            parts.append(
                f"## Nguồn: {r['source']}\n"
                f"> Kích thước: {r['chars']:,} ký tự\n\n"
                f"{r['content']}\n\n---\n"
            )
    return "\n".join(parts)


def main():
    parser = argparse.ArgumentParser(description="Gather content from files and URLs")
    parser.add_argument("--sources", help="JSON file with {files:[], urls:[]}")
    parser.add_argument("--files", nargs="*", default=[], help="Local file paths")
    parser.add_argument("--urls", nargs="*", default=[], help="URLs to fetch")
    parser.add_argument("--output", required=True, help="Output Markdown file path")
    args = parser.parse_args()

    files = list(args.files)
    urls = list(args.urls)

    # Load from JSON sources file if provided
    if args.sources:
        with open(args.sources, "r", encoding="utf-8") as f:
            src = json.load(f)
        files.extend(src.get("files", []))
        urls.extend(src.get("urls", []))

    if not files and not urls:
        print("Error: No sources provided. Use --files, --urls, or --sources.", file=sys.stderr)
        sys.exit(1)

    results = []
    success = 0
    errors = 0

    # Read local files
    for fp in files:
        r = read_local_file(fp)
        results.append(r)
        if r["error"]:
            errors += 1
        else:
            success += 1

    # Fetch URLs (with brief pause between requests to same domain)
    prev_domain = None
    for url in urls:
        from urllib.parse import urlparse
        domain = urlparse(url).netloc
        if domain == prev_domain:
            time.sleep(1)
        prev_domain = domain

        r = fetch_url(url)
        results.append(r)
        if r["error"]:
            errors += 1
        else:
            success += 1

    # Combine and write
    combined = combine_results(results)
    total_chars = sum(r["chars"] for r in results)

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(combined, encoding="utf-8")

    file_size = output_path.stat().st_size / 1024
    print(f"✅ Gathered: {output_path} ({file_size:.1f} KB, "
          f"{len(results)} sources, {success} OK, {errors} errors, "
          f"~{total_chars:,} chars)")


if __name__ == "__main__":
    main()
