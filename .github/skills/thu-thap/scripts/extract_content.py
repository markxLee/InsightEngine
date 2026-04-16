#!/usr/bin/env python3
"""
extract_content.py — Batch file content extractor for thu-thap skill.

Reads one or more local files and outputs their content as Markdown text.
Primary reader: markitdown. Falls back to format-specific libraries on failure.

Usage:
  python3 extract_content.py file1.pdf file2.docx file3.xlsx
  python3 extract_content.py --list files.txt          # read paths from a file
  python3 extract_content.py --output combined.md file1.pdf file2.docx

Output:
  Prints each file's Markdown content with a section header.
  Prints summary to stderr: "✅ Extracted: N files, X chars total"
"""

import sys
import os
import argparse
from pathlib import Path

# ---------------------------------------------------------------------------
# Primary reader: markitdown
# ---------------------------------------------------------------------------

def read_with_markitdown(path: str) -> str | None:
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(path)
        text = result.text_content.strip() if result.text_content else ""
        return text if text else None
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Format-specific fallbacks
# ---------------------------------------------------------------------------

def read_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def read_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"## Sheet: {sheet}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def read_pdf(path: str) -> str:
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
        return "\n\n".join(pages)
    except ImportError:
        pass
    # reportlab/pypdf fallback
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n\n".join(
        p.extract_text() for p in reader.pages if p.extract_text()
    )


def read_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if texts:
            slides.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def read_txt(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()


FALLBACKS = {
    ".docx": read_docx,
    ".doc":  read_docx,
    ".xlsx": read_xlsx,
    ".xlsm": read_xlsx,
    ".csv":  lambda p: open(p, encoding="utf-8", errors="replace").read(),
    ".pdf":  read_pdf,
    ".pptx": read_pptx,
    ".txt":  read_txt,
    ".md":   read_txt,
    ".html": read_txt,
}


def extract(path: str) -> tuple[str, str | None]:
    """Returns (content, error_message). error_message is None on success."""
    if not Path(path).exists():
        return "", f"File not found: {path}"

    ext = Path(path).suffix.lower()

    # Try markitdown first
    content = read_with_markitdown(path)
    if content:
        return content, None

    # Fallback by extension
    fallback_fn = FALLBACKS.get(ext)
    if fallback_fn:
        try:
            content = fallback_fn(path).strip()
            if content:
                return content, None
        except Exception as e:
            return "", f"Fallback failed for {path}: {e}"

    return "", f"Unsupported format or empty file: {path}"


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Extract content from local files to Markdown")
    parser.add_argument("files", nargs="*", help="File paths to extract")
    parser.add_argument("--list", metavar="FILE", help="Read file paths from a text file (one per line)")
    parser.add_argument("--output", metavar="FILE", help="Write combined output to this file (default: stdout)")
    args = parser.parse_args()

    paths = list(args.files)
    if args.list:
        with open(args.list, encoding="utf-8") as f:
            paths += [line.strip() for line in f if line.strip()]

    if not paths:
        parser.print_help()
        sys.exit(1)

    sections = []
    total_chars = 0
    errors = []

    for path in paths:
        content, err = extract(path)
        if err:
            errors.append(err)
            print(f"⚠️  {err}", file=sys.stderr)
        else:
            total_chars += len(content)
            sections.append(f"---\n## Source: {path}\n\n{content}")

    combined = "\n\n".join(sections)

    if args.output:
        Path(args.output).write_text(combined, encoding="utf-8")
        size_kb = round(Path(args.output).stat().st_size / 1024, 1)
        print(f"✅ Extracted: {len(sections)} files, {total_chars:,} chars → {args.output} ({size_kb} KB)", file=sys.stderr)
    else:
        print(combined)
        print(f"\n✅ Extracted: {len(sections)}/{len(paths)} files, {total_chars:,} chars total", file=sys.stderr)

    if errors:
        print(f"⚠️  {len(errors)} error(s) — see above", file=sys.stderr)
        sys.exit(1 if len(sections) == 0 else 0)


if __name__ == "__main__":
    main()
