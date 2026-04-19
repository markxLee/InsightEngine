#!/usr/bin/env python3
"""
Generate Excel (per-province sheets), a simple PPTX (summary slides), and a reveal.js HTML from extracted job JSONs.
Reads: tmp/jobs_extracted_playwright.json and tmp/jobs_extracted.json
Writes: output/jobs_by_province.xlsx, output/jobs_summary.pptx, output/reveal_jobs.html
"""
import json, os, re
from pathlib import Path

try:
    import pandas as pd
except Exception:
    pd = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

OUT_DIR = Path('output')
OUT_DIR.mkdir(exist_ok=True)

def load_jobs():
    jobs = []
    # Prefer filtered file if present
    candidate_files = []
    if os.path.exists('tmp/jobs_extracted_filtered.json'):
        candidate_files.append('tmp/jobs_extracted_filtered.json')
    candidate_files.extend(['tmp/jobs_extracted_playwright.json','tmp/jobs_extracted.json'])
    for p in candidate_files:
        if os.path.exists(p):
            try:
                data = json.load(open(p,'r',encoding='utf-8'))
                jobs.extend(data.get('jobs',[]))
            except Exception:
                pass
    # dedupe by url
    seen=set(); uniq=[]
    for j in jobs:
        url = j.get('job_detail_url') or j.get('sources',[{}])[0].get('url')
        if not url:
            url = j.get('job_detail_url','')
        if url in seen:
            continue
        seen.add(url)
        uniq.append(j)
    return uniq

def province_key(j):
    pc = j.get('province_city')
    if isinstance(pc, list) and pc:
        return pc[0]
    # try address
    addr = (j.get('company_address') or '')
    # crude match for common provinces
    provinces = ['Hà Nội','Ha Noi','Hồ Chí Minh','Ho Chi Minh','HCM','Đà Nẵng','Da Nang']
    for p in provinces:
        if p.lower() in addr.lower():
            return p
    return 'Unknown'

def sanitize_sheet(name):
    name = re.sub(r'[\\/:*?\[\]]','_', name)
    return name[:31]

def make_excel(jobs):
    if pd is None:
        print('pandas not installed; skipping Excel')
        return None
    rows = []
    for j in jobs:
        rows.append({
            'title': j.get('job_title',''),
            'company': j.get('company_name',''),
            'address': j.get('company_address',''),
            'province': province_key(j),
            'salary_raw': j.get('salary_raw',''),
            'experience_raw': j.get('years_of_experience_raw',''),
            'skills': ', '.join(j.get('required_skills',[])) if j.get('required_skills') else '',
            'url': j.get('job_detail_url') or (j.get('sources',[{}])[0].get('url') if j.get('sources') else '')
        })
    df = pd.DataFrame(rows)
    out = OUT_DIR / 'jobs_by_province.xlsx'
    with pd.ExcelWriter(out, engine='openpyxl') as w:
        for prov, group in df.groupby('province'):
            sheet = sanitize_sheet(prov or 'Unknown')
            group.to_excel(w, sheet_name=sheet, index=False)
    print('Wrote', out)
    return str(out)

def make_pptx(jobs):
    if Presentation is None:
        print('python-pptx not installed; skipping PPTX')
        return None
    prs = Presentation()
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Job Listings — Fresher/Junior (Vietnam)'
    subtitle.text = 'Generated from extracted data'
    # summary slide
    counts = {}
    for j in jobs:
        k = province_key(j) or 'Unknown'
        counts[k] = counts.get(k,0)+1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Counts by Province'
    body = slide.shapes.placeholders[1].text_frame
    for k,v in sorted(counts.items(), key=lambda x:-x[1]):
        p = body.add_paragraph()
        p.text = f'{k}: {v}'
        p.level = 1
    # add sample job slides (up to 10)
    for j in jobs[:10]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = j.get('job_title','')[:60]
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        lines = [f"Company: {j.get('company_name','')}", f"Location: {j.get('company_address','')}", f"Skills: {', '.join(j.get('required_skills',[]))}", f"URL: {j.get('job_detail_url','')}"]
        for L in lines:
            p = tf.add_paragraph()
            p.text = L
            p.level = 0
    out = OUT_DIR / 'jobs_summary.pptx'
    prs.save(out)
    print('Wrote', out)
    return str(out)

def make_reveal(jobs):
    html_lines = [
        '<!doctype html>',
        '<html>',
        '<head>',
        '  <meta charset="utf-8">',
        '  <meta name="viewport" content="width=device-width, initial-scale=1.0">',
        '  <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/reveal.css">',
        '  <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/theme/black.css">',
        '</head>',
        '<body>',
        '<div class="reveal"><div class="slides">'
    ]
    # title
    html_lines.append('<section><h1>Job Listings — Fresher/Junior (Vietnam)</h1><p>Generated</p></section>')
    # one section per province
    groups = {}
    for j in jobs:
        k = province_key(j)
        groups.setdefault(k,[]).append(j)
    for prov, items in groups.items():
        html_lines.append(f'<section><h2>{prov} — {len(items)} jobs</h2>')
        for it in items[:20]:
            title = it.get('job_title','')
            company = it.get('company_name','')
            url = it.get('job_detail_url','')
            html_lines.append(f'<article><h3>{title}</h3><p>{company}</p><p><a href="{url}">{url}</a></p></article>')
        html_lines.append('</section>')
    html_lines.extend(['</div></div>','<script src="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/reveal.js"></script>','<script>Reveal.initialize();</script>','</body>','</html>'])
    out = OUT_DIR / 'reveal_jobs.html'
    open(out,'w',encoding='utf-8').write('\n'.join(html_lines))
    print('Wrote', out)
    return str(out)

def main():
    jobs = load_jobs()
    print('Loaded', len(jobs), 'unique jobs')
    excel = make_excel(jobs)
    pptx = make_pptx(jobs)
    html = make_reveal(jobs)
    print('Done. Excel:', excel, 'PPTX:', pptx, 'HTML:', html)

if __name__ == '__main__':
    main()
