#!/usr/bin/env python3
"""
create_placeholder.py — Template-first output: create structural placeholder files.

Usage:
  create_placeholder.py <format> <output_path> [options]

Formats:
  excel    Create .xlsx with required sheets and column headers (no data)
  word     Create .docx with section headings (no content)
  slide    Create .pptx with slide stubs (title only, no content)
  html     Create .html with section scaffolding (no content)

Options:
  --sheets <comma-list>     Excel: sheet names to create
  --columns <json>          Excel: {sheet_name: [col1, col2, ...]} mapping
  --sections <comma-list>   Word/HTML: section/heading names
  --slide-titles <json>     Slide: list of slide titles
  --requirements <json>     Structured requirements from save_state.py check-requirements
  --fill <json_file>        Fill mode (US-13.4.3): fill validated placeholder with real content

Fill Mode (US-13.4.3 — update, NOT create):
  Fills an existing validated placeholder with real content.
  Never recreates the file — preserves validated structure.

  Fill JSON formats:
    Excel:  {"SheetName": [{"col1": "val1", "col2": "val2"}, ...]}
    Word:   {"title": "Title", "sections": {"Heading Name": "Content..."}}
    HTML:   {"title": "Title", "sections": {"section-1": "<p>Content</p>"}}

Example (create placeholder):
  python3 scripts/create_placeholder.py excel output/report.xlsx \\
    --sheets "Ha Noi,TP HCM" \\
    --columns '{"Ha Noi": ["STT","Company","Role","Salary"]}'

Example (fill placeholder — US-13.4.3):
  python3 scripts/create_placeholder.py excel output/report.xlsx \\
    --fill tmp/excel_content.json
  python3 scripts/create_placeholder.py word output/report.docx \\
    --fill tmp/word_content.json
"""

import sys
import json
from pathlib import Path

SCRIPTS_DIR = Path(__file__).parent
OUTPUT_DIR = SCRIPTS_DIR.parent / "output"
TMP_DIR = SCRIPTS_DIR.parent / "tmp"


def parse_args():
    args = sys.argv[1:]
    if len(args) < 2:
        print("Usage: create_placeholder.py <format> <output_path> [options]")
        print("Formats: excel, word, slide, html")
        sys.exit(1)

    fmt = args[0].lower()
    output_path = Path(args[1])
    options = {}

    i = 2
    while i < len(args):
        if args[i] == "--sheets" and i + 1 < len(args):
            options["sheets"] = [s.strip() for s in args[i + 1].split(",")]
            i += 2
        elif args[i] == "--columns" and i + 1 < len(args):
            options["columns"] = json.loads(args[i + 1])
            i += 2
        elif args[i] == "--sections" and i + 1 < len(args):
            options["sections"] = [s.strip() for s in args[i + 1].split(",")]
            i += 2
        elif args[i] == "--slide-titles" and i + 1 < len(args):
            options["slide_titles"] = json.loads(args[i + 1])
            i += 2
        elif args[i] == "--fill" and i + 1 < len(args):
            options["fill"] = Path(args[i + 1])
            i += 2
        elif args[i] == "--requirements" and i + 1 < len(args):
            try:
                options["requirements"] = json.loads(args[i + 1])
            except json.JSONDecodeError:
                pass
            i += 2
        else:
            i += 1

    return fmt, output_path, options


def infer_from_requirements(requirements: dict) -> dict:
    """Infer placeholder structure from structured_requirements JSON."""
    result = {}

    # Grouping → sheets for Excel
    grouping = requirements.get("grouping", [])
    if grouping:
        result["sheets"] = grouping

    # fields_required → columns (applied to all sheets if not per-sheet)
    fields = requirements.get("fields_required", {})
    if fields:
        if isinstance(fields, dict):
            result["columns"] = fields  # {sheet_name: [col1, col2]}
        elif isinstance(fields, list):
            result["default_columns"] = fields

    # content_requirements → sections for Word/HTML
    content_reqs = requirements.get("content_requirements", [])
    if content_reqs:
        result["sections"] = content_reqs

    return result


# ─────────────────────────────────────
# Excel placeholder
# ─────────────────────────────────────

def create_excel_placeholder(output_path: Path, options: dict):
    try:
        import openpyxl
    except ImportError:
        print("Error: openpyxl not installed. Run: pip install openpyxl")
        sys.exit(1)

    # Infer structure from requirements if provided
    if "requirements" in options:
        inferred = infer_from_requirements(options["requirements"])
        for k, v in inferred.items():
            if k not in options:
                options[k] = v

    sheets = options.get("sheets", ["Sheet1"])
    columns_map = options.get("columns", {})
    default_columns = options.get("default_columns", ["STT", "Tên cột 1", "Tên cột 2"])

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # Remove default sheet

    for sheet_name in sheets:
        ws = wb.create_sheet(title=sheet_name[:31])  # Excel max 31 chars
        cols = columns_map.get(sheet_name, default_columns)
        for col_idx, col_name in enumerate(cols, start=1):
            ws.cell(row=1, column=col_idx, value=col_name)
        # Style header row
        from openpyxl.styles import Font, PatternFill, Alignment
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")
            cell.alignment = Alignment(horizontal="center")
        # Add placeholder data row (visual indicator, not real data)
        ws.append(["[PLACEHOLDER — fill with real data]"] + [""] * (len(cols) - 1))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(output_path)
    size = output_path.stat().st_size
    sheet_list = ", ".join(sheets[:5]) + ("..." if len(sheets) > 5 else "")
    print(f"PLACEHOLDER_CREATED: {output_path} ({size} bytes, {len(sheets)} sheets: {sheet_list})")


# ─────────────────────────────────────
# Word placeholder
# ─────────────────────────────────────

def create_word_placeholder(output_path: Path, options: dict):
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        print("Error: python-docx not installed. Run: pip install python-docx")
        sys.exit(1)

    if "requirements" in options:
        inferred = infer_from_requirements(options["requirements"])
        for k, v in inferred.items():
            if k not in options:
                options[k] = v

    sections = options.get("sections", ["Phần 1", "Phần 2", "Phần 3"])

    doc = Document()
    doc.add_paragraph("[PLACEHOLDER DOCUMENT — fill sections with real content]")
    for i, section in enumerate(sections, start=1):
        doc.add_heading(section, level=1)
        p = doc.add_paragraph(f"[PLACEHOLDER — nội dung phần {i}: {section}]")
        p.runs[0].italic = True

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    size = output_path.stat().st_size
    print(f"PLACEHOLDER_CREATED: {output_path} ({size} bytes, {len(sections)} sections)")


# ─────────────────────────────────────
# PowerPoint placeholder
# ─────────────────────────────────────

def create_slide_placeholder(output_path: Path, options: dict):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)

    if "requirements" in options:
        inferred = infer_from_requirements(options["requirements"])
        sections = inferred.get("sections", [])
        if sections and "slide_titles" not in options:
            options["slide_titles"] = sections

    titles = options.get("slide_titles", [f"Slide {i}" for i in range(1, 6)])

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # blank

    for title_text in titles:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        txBox.text_frame.text = f"[PLACEHOLDER] {title_text}"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    size = output_path.stat().st_size
    print(f"PLACEHOLDER_CREATED: {output_path} ({size} bytes, {len(titles)} slides)")


# ─────────────────────────────────────
# HTML placeholder
# ─────────────────────────────────────

def create_html_placeholder(output_path: Path, options: dict):
    if "requirements" in options:
        inferred = infer_from_requirements(options["requirements"])
        for k, v in inferred.items():
            if k not in options:
                options[k] = v

    sections = options.get("sections", ["Phần 1", "Phần 2", "Phần 3"])

    section_html = "\n".join(
        f'  <section id="section-{i}">\n'
        f'    <h2>{s}</h2>\n'
        f'    <p class="placeholder">[PLACEHOLDER — nội dung: {s}]</p>\n'
        f'  </section>'
        for i, s in enumerate(sections, start=1)
    )

    html = f"""<!DOCTYPE html>
<html lang="vi">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>[PLACEHOLDER]</title>
  <style>
    body {{ font-family: sans-serif; max-width: 900px; margin: 0 auto; padding: 2rem; }}
    .placeholder {{ color: #999; font-style: italic; background: #f5f5f5; padding: 0.5rem; }}
    h2 {{ border-bottom: 1px solid #ddd; padding-bottom: 0.3rem; }}
  </style>
</head>
<body>
  <h1>[PLACEHOLDER DOCUMENT]</h1>
{section_html}
</body>
</html>
"""

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(html, encoding="utf-8")
    size = output_path.stat().st_size
    print(f"PLACEHOLDER_CREATED: {output_path} ({size} bytes, {len(sections)} sections)")


# ─────────────────────────────────────
# Fill mode
# ─────────────────────────────────────

def fill_excel_placeholder(output_path: Path, fill_path: Path):
    """Fill an existing Excel placeholder with data from a JSON file.
    
    JSON format:
    {
      "SheetName": [
        {"col1": "val1", "col2": "val2"},
        ...
      ]
    }
    """
    try:
        import openpyxl
    except ImportError:
        print("Error: openpyxl not installed"); sys.exit(1)

    if not output_path.exists():
        print(f"Error: placeholder file not found: {output_path}"); sys.exit(1)

    content = json.loads(fill_path.read_text(encoding="utf-8"))
    wb = openpyxl.load_workbook(output_path)

    for sheet_name, rows in content.items():
        if sheet_name not in wb.sheetnames:
            print(f"Warning: sheet '{sheet_name}' not in placeholder — skipping")
            continue
        ws = wb[sheet_name]
        # Get column order from header row
        headers = [ws.cell(row=1, column=c).value for c in range(1, ws.max_column + 1)]
        # Clear placeholder data row (row 2)
        if ws.max_row >= 2 and ws.cell(row=2, column=1).value and str(ws.cell(row=2, column=1).value).startswith("[PLACEHOLDER"):
            ws.delete_rows(2)
        # Append real data
        for row_data in rows:
            row = [row_data.get(h, "") for h in headers]
            ws.append(row)

    wb.save(output_path)
    size = output_path.stat().st_size
    print(f"PLACEHOLDER_FILLED: {output_path} ({size} bytes)")


def fill_word_placeholder(output_path: Path, fill_path: Path):
    """Fill an existing Word placeholder with content from a JSON file.
    
    JSON format:
    {
      "title": "Document Title",
      "sections": {
        "Section Name": "Content for this section...",
        "Another Section": "More content..."
      }
    }
    """
    try:
        from docx import Document
    except ImportError:
        print("Error: python-docx not installed"); sys.exit(1)

    if not output_path.exists():
        print(f"Error: placeholder file not found: {output_path}"); sys.exit(1)

    content = json.loads(fill_path.read_text(encoding="utf-8"))
    doc = Document(output_path)

    section_content = content.get("sections", {})
    title = content.get("title", "")

    # Update title if provided
    if title and doc.paragraphs:
        for para in doc.paragraphs:
            if para.style.name == "Normal" and "[PLACEHOLDER" in para.text:
                para.clear()
                para.add_run(title)
                break

    # Update sections: find heading, replace placeholder paragraph below it
    i = 0
    while i < len(doc.paragraphs):
        para = doc.paragraphs[i]
        if para.style.name.startswith("Heading") and para.text in section_content:
            section_text = section_content[para.text]
            # Find placeholder paragraph after this heading
            if i + 1 < len(doc.paragraphs):
                next_para = doc.paragraphs[i + 1]
                if "[PLACEHOLDER" in next_para.text:
                    next_para.clear()
                    # Split content by newlines and add as separate runs
                    for line in section_text.split("\n"):
                        next_para.add_run(line + " ")
        i += 1

    doc.save(output_path)
    size = output_path.stat().st_size
    print(f"PLACEHOLDER_FILLED: {output_path} ({size} bytes)")


def fill_html_placeholder(output_path: Path, fill_path: Path):
    """Fill an existing HTML placeholder with content from a JSON file.
    
    JSON format:
    {
      "title": "Page Title",
      "sections": {
        "section-1": "<p>Content for section 1</p>",
        "section-2": "<p>Content for section 2</p>"
      }
    }
    The section keys are matched against section id attributes in the HTML.
    """
    if not output_path.exists():
        print(f"Error: placeholder file not found: {output_path}"); sys.exit(1)

    content = json.loads(fill_path.read_text(encoding="utf-8"))
    html = output_path.read_text(encoding="utf-8")

    # Update title
    if content.get("title"):
        html = html.replace("[PLACEHOLDER DOCUMENT]", content["title"])
        html = html.replace("<title>[PLACEHOLDER]</title>", f"<title>{content['title']}</title>")

    # Update sections — replace placeholder paragraph inside matching section
    sections = content.get("sections", {})
    for section_id, section_html in sections.items():
        # Match section by id or heading text
        import re
        # Replace placeholder paragraph in section with real content
        pattern = rf'(<section[^>]*id="{section_id}"[^>]*>.*?<p class="placeholder">)[^<]+(</p>)'
        replacement = rf'\g<1>{section_html}\2'
        new_html = re.sub(pattern, replacement, html, flags=re.DOTALL)
        if new_html != html:
            html = new_html

    output_path.write_text(html, encoding="utf-8")
    size = output_path.stat().st_size
    print(f"PLACEHOLDER_FILLED: {output_path} ({size} bytes)")


# ─────────────────────────────────────
# Main
# ─────────────────────────────────────

def main():
    fmt, output_path, options = parse_args()

    if "fill" in options:
        fill_path = options["fill"]
        if fmt == "excel":
            fill_excel_placeholder(output_path, fill_path)
        elif fmt in ("word", "docx"):
            fill_word_placeholder(output_path, fill_path)
        elif fmt in ("html",):
            fill_html_placeholder(output_path, fill_path)
        else:
            print(f"Fill mode not yet implemented for format: {fmt}")
            sys.exit(1)
        return

    if fmt == "excel":
        create_excel_placeholder(output_path, options)
    elif fmt in ("word", "docx"):
        create_word_placeholder(output_path, options)
    elif fmt in ("slide", "pptx", "ppt"):
        create_slide_placeholder(output_path, options)
    elif fmt in ("html",):
        create_html_placeholder(output_path, options)
    else:
        print(f"Error: unknown format '{fmt}'. Supported: excel, word, slide, html")
        sys.exit(1)


if __name__ == "__main__":
    main()
